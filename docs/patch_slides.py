"""Patch slide pptx — set background, add logo, replace fake AI images, fix text."""
from copy import deepcopy
from pptx import Presentation
from pptx.util import Emu, Inches
from pptx.oxml.ns import qn
from lxml import etree

SRC = 'Phan-Mem-Mua-Ban-Trao-DJoi-DJien-Thoai-Tich-Hop-AI-Ho-Tro-DJinh-Gia-San-Pham.pptx'
OUT = SRC
BG = 'slide-background.png'
LOGO = 'logo-duytan.png'
VAL_BATCH = 'damage-model-artifacts/val_batch0_pred.jpg'
CONF_MATRIX = 'damage-model-artifacts/confusion_matrix.png'

# Banner AI trang trí cần xóa — dựa trên inventory:
# slide 1 -> Image 0 (3936KB) — right banner
# slide 3 -> Image 0 (6301KB) — right banner
# slide 5 -> Image 0 (5050KB) — left banner
# slide 10 -> Image 0 (5302KB) — left banner
# slide 15 -> Image 0 (4352KB) — left banner
# slide 17 -> Image 0 (5721KB) — right banner
# slide 18 -> Image 0 (1499KB) — left banner (smaller)
SLIDES_REMOVE_BANNER = {1, 3, 5, 10, 15, 17, 18}
BANNER_MIN_BYTES = 1_300_000  # ~1.3 MB — > BG (1066KB) nhưng < val_batch (~550KB sau)

# Slide 12 — thay ảnh AI minh họa (~4.4 MB) bằng val_batch0_pred
SLIDE_REPLACE_YOLO = 12

# Slide 16 — minh họa chat AI (~3.1 MB) — XÓA (chưa thay vì user phải tự chụp)
SLIDE_REMOVE_CHAT_ILLUSTRATION = 16

# Slide 14 — 2 mockup AI giả ~1.1 MB mỗi cái — XÓA (user tự chụp screenshot UI)
SLIDE_REMOVE_HOMEPAGE_MOCKUP = 14

prs = Presentation(SRC)
slide_w = prs.slide_width
slide_h = prs.slide_height
print(f"Slides: {len(prs.slides)} | Size: {slide_w}x{slide_h} EMU")


def set_slide_background(slide, image_path):
    """Set background bằng cách insert picture phủ toàn bộ + đẩy xuống dưới cùng z-order."""
    pic = slide.shapes.add_picture(image_path, 0, 0, width=slide_w, height=slide_h)
    # Đẩy xuống đáy z-order
    sp_tree = pic._element.getparent()
    sp_tree.remove(pic._element)
    # Insert ngay sau <p:nvGrpSpPr> + <p:grpSpPr> (2 elem đầu)
    # spTree structure: nvGrpSpPr (1) + grpSpPr (2) + shapes...
    # Insert tại index 2 (vị trí thứ 3, ngay sau 2 elem header)
    sp_tree.insert(2, pic._element)
    return pic


# === STEP 1: Set background cho tất cả slide ===
for i, slide in enumerate(prs.slides, 1):
    set_slide_background(slide, BG)
    print(f"[BG] Slide {i}: added background")

# === STEP 2: Xóa banner AI trang trí ở các slide đã chọn ===
for i, slide in enumerate(prs.slides, 1):
    if i not in SLIDES_REMOVE_BANNER:
        continue
    to_remove = []
    for shape in slide.shapes:
        if shape.shape_type == 13:  # Picture
            try:
                blob = shape.image.blob
                if len(blob) >= BANNER_MIN_BYTES:
                    to_remove.append(shape)
            except Exception:
                pass
    for shape in to_remove:
        sp = shape._element
        sp.getparent().remove(sp)
        print(f"[RM-BANNER] Slide {i}: removed picture {len(blob)//1024}KB")

# === STEP 3: Slide 12 — thay YOLO illustration bằng val_batch0_pred.jpg ===
slide12 = prs.slides[SLIDE_REPLACE_YOLO - 1]
yolo_pic = None
for shape in slide12.shapes:
    if shape.shape_type == 13:
        try:
            if len(shape.image.blob) > 2 * 1024 * 1024:  # 4.4MB image
                yolo_pic = shape
                break
        except Exception:
            pass

if yolo_pic is not None:
    left, top, width, height = yolo_pic.left, yolo_pic.top, yolo_pic.width, yolo_pic.height
    sp = yolo_pic._element
    sp.getparent().remove(sp)
    new_pic = slide12.shapes.add_picture(VAL_BATCH, left, top, width=width, height=height)
    print(f"[YOLO] Slide 12: replaced AI illustration with val_batch0_pred.jpg")

# === STEP 4: Slide 14 — xóa 2 mockup AI giả (~1.1MB mỗi cái)
# CHỈ xóa picture KHÔNG phải BG (BG được add đầu tiên, ở z-bottom)
slide14 = prs.slides[SLIDE_REMOVE_HOMEPAGE_MOCKUP - 1]
to_remove = []
bg_seen = False
for shape in slide14.shapes:
    if shape.shape_type == 13:
        try:
            size = len(shape.image.blob)
            # Skip BG (first picture, ~1066KB) — chỉ xóa mockup nhỏ hơn BG khoảng 1.1MB
            if not bg_seen:
                bg_seen = True
                continue
            if 800 * 1024 < size < 1_300_000:
                to_remove.append(shape)
        except Exception:
            pass
for shape in to_remove:
    sp = shape._element
    sp.getparent().remove(sp)
    print(f"[RM-MOCKUP] Slide 14: removed fake mockup")

# === STEP 5: Slide 16 — xóa minh họa chat AI ===
slide16 = prs.slides[SLIDE_REMOVE_CHAT_ILLUSTRATION - 1]
to_remove = []
for shape in slide16.shapes:
    if shape.shape_type == 13:
        try:
            if len(shape.image.blob) > 2 * 1024 * 1024:
                to_remove.append(shape)
        except Exception:
            pass
for shape in to_remove:
    sp = shape._element
    sp.getparent().remove(sp)
    print(f"[RM-CHATILL] Slide 16: removed AI chat illustration")

# === STEP 6: Slide 18 — chèn confusion_matrix.png góc dưới phải ===
slide18 = prs.slides[18 - 1]
# Đặt ở góc dưới phải: left=55% top=55%, width=42% height=42%
cm_left = int(slide_w * 0.55)
cm_top = int(slide_h * 0.50)
cm_w = int(slide_w * 0.42)
cm_h = int(slide_h * 0.45)
slide18.shapes.add_picture(CONF_MATRIX, cm_left, cm_top, width=cm_w, height=cm_h)
print(f"[CONF] Slide 18: added confusion_matrix.png")

# === STEP 7: Slide 1 — chèn logo DTU góc trên trái ===
slide1 = prs.slides[0]
logo_left = int(slide_w * 0.04)
logo_top = int(slide_h * 0.05)
logo_w = int(slide_w * 0.08)  # 8% width
logo_h = int(slide_h * 0.16)  # 16% height (theo tỷ lệ logo ~1:1)
slide1.shapes.add_picture(LOGO, logo_left, logo_top, width=logo_w, height=logo_h)
print(f"[LOGO] Slide 1: added DTU logo")

# === STEP 8: Text fixes ===
def replace_text_in_slide(slide, old, new):
    found = False
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if old in run.text:
                        run.text = run.text.replace(old, new)
                        found = True
    return found

# Slide 9: 40 ngày → 6 tuần
if replace_text_in_slide(prs.slides[8], '40 ngày', '6 tuần'):
    print('[TEXT] Slide 9: 40 ngày → 6 tuần')

# Slide 17: bổ sung "35 testcase AI" — cần thêm KPI mới
# Cách đơn giản: tìm shape có text "17" và "Unit Tests", thêm note "+ 35 testcase AI"
# Hoặc thay text "17/17 pass" → "17/17 pass · 35 testcase AI"
if replace_text_in_slide(prs.slides[16], '17/17 pass', '17/17 + 35 testcase AI'):
    print('[TEXT] Slide 17: 17/17 pass → 17/17 + 35 testcase AI')

prs.save(OUT)
print(f'\nDONE — saved {OUT}')
