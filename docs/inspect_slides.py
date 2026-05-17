"""Inspect layout của slide 1 và slide 17 trước khi patch."""
from pptx import Presentation
from pptx.util import Emu

SRC = 'Phan-Mem-Mua-Ban-Trao-DJoi-DJien-Thoai-Tich-Hop-AI-Ho-Tro-DJinh-Gia-San-Pham.pptx'
prs = Presentation(SRC)
sw, sh = prs.slide_width, prs.slide_height
print(f"Slide size: {sw} x {sh} EMU = {Emu(sw).inches:.2f} x {Emu(sh).inches:.2f} inch\n")

def fmt(emu):
    return f"{Emu(emu).inches:.2f}in"

def fmt_pct(emu, total):
    return f"{emu/total*100:.1f}%"

for slide_idx in [0, 16]:  # slide 1 và 17 (0-indexed)
    slide = prs.slides[slide_idx]
    print(f"\n{'='*70}")
    print(f"SLIDE {slide_idx + 1}")
    print(f"{'='*70}")
    for i, shape in enumerate(slide.shapes):
        try:
            l, t, w, h = shape.left, shape.top, shape.width, shape.height
            kind = shape.shape_type
            name = shape.name
            text = ''
            if shape.has_text_frame:
                text = shape.text_frame.text[:80].replace('\n', ' | ')
            print(f"  [{i}] type={kind} name='{name}'")
            print(f"      pos: L={fmt(l)}({fmt_pct(l,sw)}) T={fmt(t)}({fmt_pct(t,sh)})")
            print(f"      size: W={fmt(w)}({fmt_pct(w,sw)}) H={fmt(h)}({fmt_pct(h,sh)})")
            if text:
                print(f"      text: {text!r}")
        except Exception as e:
            print(f"  [{i}] ERROR: {e}")
