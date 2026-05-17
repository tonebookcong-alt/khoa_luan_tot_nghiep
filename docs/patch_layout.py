"""Patch layout slide 1 (tinh chỉnh spacing) và slide 17 (restructure 2x2 + KPI row)."""
from pptx import Presentation
from pptx.util import Inches, Emu

SRC = 'Phan-Mem-Mua-Ban-Trao-DJoi-DJien-Thoai-Tich-Hop-AI-Ho-Tro-DJinh-Gia-San-Pham.pptx'
OUT = SRC

prs = Presentation(SRC)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def set_pos(shape, left=None, top=None, width=None, height=None):
    """Set absolute position/size (inches) for a shape."""
    if left is not None:
        shape.left = Inches(left)
    if top is not None:
        shape.top = Inches(top)
    if width is not None:
        shape.width = Inches(width)
    if height is not None:
        shape.height = Inches(height)


# =====================================================================
# SLIDE 1 — tinh chỉnh spacing (giãn đều, dịch xuống một chút)
# =====================================================================
s1 = prs.slides[0]
shapes1 = list(s1.shapes)

# [1] Badge BG, [2] Badge text — đẩy hơi xuống cho cân
set_pos(shapes1[1], left=0.54, top=0.65, width=2.84, height=0.32)
set_pos(shapes1[2], left=0.64, top=0.69, width=2.66, height=0.24)

# [3] Title — mở rộng width + đẩy lên cho cân
set_pos(shapes1[3], left=0.54, top=1.20, width=7.40, height=2.00)

# [4] University name — giãn xuống
set_pos(shapes1[4], left=0.54, top=3.45, width=8.00, height=0.30)

# [5] GVHD + thành viên — giãn xuống
set_pos(shapes1[5], left=0.54, top=3.95, width=8.00, height=0.80)

# [6] Date — đẩy xuống đáy
set_pos(shapes1[6], left=0.54, top=4.95, width=8.00, height=0.30)

# [7] DTU logo — phóng to nhẹ, đẩy ra mép phải hơn
set_pos(shapes1[7], left=8.80, top=0.30, width=0.95, height=1.05)

print("[SLIDE 1] Layout tinh chỉnh xong.")


# =====================================================================
# SLIDE 17 — restructure: KPI 4-cột ngang + Cards 2x2 grid
# =====================================================================
s17 = prs.slides[16]
shapes17 = list(s17.shapes)

# Title block (Shape 0, Text 1, Text 2) — giữ nguyên vị trí, chỉ mở rộng
# [1] KẾT QUẢ badge — giữ nguyên
# [2] Text "KẾT QUẢ" — giữ nguyên
# [3] Title "Những Công Việc Đã Hoàn Thành" — mở rộng
set_pos(shapes17[3], left=0.45, top=0.45, width=9.00, height=0.50)

# --- KPI ROW: 4 cột ngang tại T=1.20 ---
# Slide width 10in, padding 0.45 mỗi bên → vùng dùng được 9.10in
# 4 cols, gap 0.15in giữa → mỗi col = (9.10 - 3*0.15) / 4 = 2.16in
COL_W = 2.16
COL_GAP = 0.15
COL_L = [0.45, 0.45 + COL_W + COL_GAP, 0.45 + 2*(COL_W + COL_GAP), 0.45 + 3*(COL_W + COL_GAP)]
# = [0.45, 2.76, 5.07, 7.38]

KPI_TOP = 1.20         # đỉnh khu KPI
NUM_H = 0.70           # height của số to
LABEL_TOP = KPI_TOP + NUM_H + 0.05    # 1.95
LABEL_H = 0.25
DESC_TOP = LABEL_TOP + LABEL_H + 0.02  # 2.22
DESC_H = 0.25

# Mỗi KPI có 3 shape: [number, label, description]
# KPI 1: shapes [4]=14, [5]=User Stories, [6]=100% backlog
# KPI 2: shapes [7]=6,  [8]=Docker Services, [9]=Đóng gói...
# KPI 3: shapes [10]=17,[11]=Unit Tests, [12]=17/17 + 35 testcase AI
# KPI 4: shapes [13]=<3s,[14]=Latency P95, [15]=Endpoint định giá AI
kpi_shape_groups = [
    (shapes17[4], shapes17[5], shapes17[6]),
    (shapes17[7], shapes17[8], shapes17[9]),
    (shapes17[10], shapes17[11], shapes17[12]),
    (shapes17[13], shapes17[14], shapes17[15]),
]

for col_idx, (num, label, desc) in enumerate(kpi_shape_groups):
    L = COL_L[col_idx]
    # Số to — chiếm cả col width, height NUM_H
    set_pos(num, left=L, top=KPI_TOP, width=COL_W, height=NUM_H)
    # Label — centered trong col
    set_pos(label, left=L, top=LABEL_TOP, width=COL_W, height=LABEL_H)
    # Description
    set_pos(desc, left=L, top=DESC_TOP, width=COL_W, height=DESC_H)

print("[SLIDE 17] KPI row 4-cột ngang xong.")

# --- CARDS 2x2 GRID tại T=2.95 ---
# 2 cards/row, gap 0.20 giữa, padding 0.45 mỗi bên
# Card W = (10 - 0.45*2 - 0.20) / 2 = 4.45
CARD_W = 4.45
CARD_H = 1.05
CARD_GAP_X = 0.20
CARD_GAP_Y = 0.20
CARD_TOP_ROW1 = 2.95
CARD_TOP_ROW2 = CARD_TOP_ROW1 + CARD_H + CARD_GAP_Y  # 4.20

CARD_POS = [
    (0.45, CARD_TOP_ROW1),                          # Card 1: top-left
    (0.45 + CARD_W + CARD_GAP_X, CARD_TOP_ROW1),    # Card 2: top-right
    (0.45, CARD_TOP_ROW2),                          # Card 3: bottom-left
    (0.45 + CARD_W + CARD_GAP_X, CARD_TOP_ROW2),    # Card 4: bottom-right
]

# Mỗi card có 5 shape (theo thứ tự): bg_rect, icon_bg, icon_pic, title_text, desc_text
# Card 1: [16][17][18][19][20]  Backend + Frontend
# Card 2: [21][22][23][24][25]  AI Core
# Card 3: [26][27][28][29][30]  Giao diện
# Card 4: [31][32][33][34][35]  Dữ liệu Polyglot
card_shape_groups = [
    (shapes17[16], shapes17[17], shapes17[18], shapes17[19], shapes17[20]),
    (shapes17[21], shapes17[22], shapes17[23], shapes17[24], shapes17[25]),
    (shapes17[26], shapes17[27], shapes17[28], shapes17[29], shapes17[30]),
    (shapes17[31], shapes17[32], shapes17[33], shapes17[34], shapes17[35]),
]

# Offset trong card (relative tới card origin):
ICON_BG_OFF = (0.20, 0.20)
ICON_BG_SIZE = (0.45, 0.45)
ICON_PIC_OFF = (0.32, 0.32)
ICON_PIC_SIZE = (0.20, 0.20)
TITLE_OFF = (0.20, 0.75)        # title nằm dưới icon
TITLE_W = 1.80
TITLE_H = 0.22
DESC_OFF = (0.20, 0.78)         # nằm cùng row với title, bên phải
# Vì card mới rộng hơn (4.45) → có thể đặt title-desc theo hàng ngang
# title: 0.85 đến ~2.30, desc: từ 2.40 đến hết
TITLE_X = 0.85
TITLE_Y = 0.30
DESC_X = 0.85
DESC_Y = 0.62

for i, (bg, icon_bg, icon_pic, title, desc) in enumerate(card_shape_groups):
    cx, cy = CARD_POS[i]
    # Card BG
    set_pos(bg, left=cx, top=cy, width=CARD_W, height=CARD_H)
    # Icon circle BG — căn giữa dọc trong card
    icon_bg_y = cy + (CARD_H - ICON_BG_SIZE[1]) / 2
    set_pos(icon_bg, left=cx + ICON_BG_OFF[0], top=icon_bg_y,
            width=ICON_BG_SIZE[0], height=ICON_BG_SIZE[1])
    # Icon picture — căn giữa trong icon BG
    icon_pic_x = cx + ICON_BG_OFF[0] + (ICON_BG_SIZE[0] - ICON_PIC_SIZE[0]) / 2
    icon_pic_y = icon_bg_y + (ICON_BG_SIZE[1] - ICON_PIC_SIZE[1]) / 2
    set_pos(icon_pic, left=icon_pic_x, top=icon_pic_y,
            width=ICON_PIC_SIZE[0], height=ICON_PIC_SIZE[1])
    # Title — bên phải icon, dòng trên
    set_pos(title, left=cx + TITLE_X, top=cy + TITLE_Y, width=TITLE_W, height=TITLE_H)
    # Description — bên phải icon, dòng dưới
    set_pos(desc, left=cx + DESC_X, top=cy + DESC_Y,
            width=CARD_W - DESC_X - 0.15, height=0.25)

print("[SLIDE 17] Cards 2x2 grid xong.")

prs.save(OUT)
print(f"\nSAVED → {OUT}")
